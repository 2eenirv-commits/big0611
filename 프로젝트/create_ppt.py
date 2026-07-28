"""부동산 웹 플랫폼 발표용 PPT 생성 스크립트 (10장, 20분)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 색상 테마 (Bootstrap 5 스타일)
PRIMARY = RGBColor(0x0D, 0x6E, 0xFD)
DARK = RGBColor(0x21, 0x25, 0x29)
GRAY = RGBColor(0x6C, 0x75, 0x7D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
ACCENT = RGBColor(0x19, 0x87, 0x54)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TITLE_LAYOUT = prs.slide_layouts[0]


def add_header_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(0.5), Inches(0.65), Inches(12), Inches(0.35))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(14)
        sp.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)


def add_bullets(slide, items, left=0.6, top=1.4, width=12, height=5.5, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_table(slide, headers, rows, left=0.6, top=1.5, col_widths=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [Inches(12 / n_cols)] * n_cols
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(12), Inches(0.5 * n_rows))
    table = table_shape.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    return table_shape


# ── 슬라이드 1: 표지 ──
slide = prs.slides.add_slide(BLANK)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.0), prs.slide_width, Inches(0.08))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = PRIMARY
accent_bar.line.fill.background()

title_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.2))
tp = title_box.text_frame.paragraphs[0]
tp.text = "부동산 웹 플랫폼"
tp.font.size = Pt(44)
tp.font.bold = True
tp.font.color.rgb = WHITE
tp.alignment = PP_ALIGN.CENTER

sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.3), Inches(11), Inches(0.8))
sp = sub_box.text_frame.paragraphs[0]
sp.text = "제작 결과 발표  |  매물 검색 · 분석 · 자동화 통합 플랫폼"
sp.font.size = Pt(22)
sp.font.color.rgb = RGBColor(0xAD, 0xD8, 0xFF)
sp.alignment = PP_ALIGN.CENTER

team_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(0.6))
team_p = team_box.text_frame.paragraphs[0]
team_p.text = "팀원: 김길동  ·  홍길동  ·  박미금  ·  김서연"
team_p.font.size = Pt(18)
team_p.font.color.rgb = GRAY
team_p.alignment = PP_ALIGN.CENTER

date_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
date_p = date_box.text_frame.paragraphs[0]
date_p.text = "빅데이터 분석 & 자동화 시스템 구축  |  2026.07.28"
date_p.font.size = Pt(16)
date_p.font.color.rgb = GRAY
date_p.alignment = PP_ALIGN.CENTER

time_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.4))
time_p = time_box.text_frame.paragraphs[0]
time_p.text = "발표 시간: 20분  |  10장"
time_p.font.size = Pt(14)
time_p.font.color.rgb = GRAY
time_p.alignment = PP_ALIGN.CENTER

# ── 슬라이드 2: 팀 소개 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "팀 소개", "Team Members")

members = [
    ("김길동", "PM / 백엔드 리드", "Spring Boot API · DB 설계 · Spring Security"),
    ("홍길동", "풀스택 개발", "프론트·백 연동 · Gradle 빌드 · 배포"),
    ("박미금", "프론트엔드", "HTML/CSS/JS · Bootstrap 5 · Flexbox UI"),
    ("김서연", "RPA · AI / 데이터", "Python 크롤링 · 시각화 · ML 분석"),
]

card_w = Inches(2.8)
gap = Inches(0.35)
start_x = Inches(0.55)
card_y = Inches(1.5)

for idx, (name, role, desc) in enumerate(members):
    x = start_x + idx * (card_w + gap)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, Inches(4.8))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.color.rgb = PRIMARY

    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), card_y + Inches(0.3), Inches(1.0), Inches(1.0))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = PRIMARY
    avatar.line.fill.background()
    init = slide.shapes.add_textbox(x + Inches(0.9), card_y + Inches(0.55), Inches(1.0), Inches(0.5))
    ip = init.text_frame.paragraphs[0]
    ip.text = name[0]
    ip.font.size = Pt(28)
    ip.font.bold = True
    ip.font.color.rgb = WHITE
    ip.alignment = PP_ALIGN.CENTER

    nb = slide.shapes.add_textbox(x + Inches(0.15), card_y + Inches(1.5), card_w - Inches(0.3), Inches(0.5))
    np = nb.text_frame.paragraphs[0]
    np.text = name
    np.font.size = Pt(20)
    np.font.bold = True
    np.font.color.rgb = DARK
    np.alignment = PP_ALIGN.CENTER

    rb = slide.shapes.add_textbox(x + Inches(0.15), card_y + Inches(2.0), card_w - Inches(0.3), Inches(0.4))
    rp = rb.text_frame.paragraphs[0]
    rp.text = role
    rp.font.size = Pt(14)
    rp.font.bold = True
    rp.font.color.rgb = PRIMARY
    rp.alignment = PP_ALIGN.CENTER

    db = slide.shapes.add_textbox(x + Inches(0.15), card_y + Inches(2.5), card_w - Inches(0.3), Inches(2.0))
    db.text_frame.word_wrap = True
    dp = db.text_frame.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = GRAY
    dp.alignment = PP_ALIGN.CENTER

collab = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(12), Inches(0.5))
cp = collab.text_frame.paragraphs[0]
cp.text = "협업: GitHub · Gradle · API 명세 공유 · 주 2회 스프린트 회의"
cp.font.size = Pt(14)
cp.font.color.rgb = GRAY
cp.alignment = PP_ALIGN.CENTER

# ── 슬라이드 3: 프로젝트 개요 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "프로젝트 개요", "Project Overview")

add_bullets(slide, [
    "배경",
    ("부동산 정보가 여러 사이트에 분산 → 비교·분석 어려움", 1),
    ("수동 검색·엑셀 정리에 시간 과다 소요", 1),
    "",
    "목표",
    ("통합 매물 검색 — 지역·가격·면적 필터", 1),
    ("회원·보안 — Spring Security + BCrypt", 1),
    ("데이터 분석 — Python 시각화·가격 예측", 1),
    ("자동 수집(RPA) — 매물 정보 크롤링", 1),
    "",
    "기대 효과: 검색 시간 단축 · 데이터 기반 의사결정 · 풀스택 역량 통합",
], font_size=17)

# ── 슬라이드 4: 시스템 아키텍처 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "시스템 아키텍처", "System Architecture")

layers = [
    ("사용자 브라우저", "HTML · CSS · JS · Bootstrap 5 · Flexbox", PRIMARY),
    ("Spring Boot 3.5 + Thymeleaf + Tomcat 10.1", "Spring Web · Security · JDBC Template", RGBColor(0x0A, 0x58, 0xCA)),
    ("MySQL 8  /  TiDB Cloud", "매물 · 회원 · 분석 데이터 저장", ACCENT),
    ("RPA + AI (Python)", "NumPy · Pandas · Matplotlib · Seaborn · Plotly · Scikit-learn", RGBColor(0xDC, 0x35, 0x45)),
]

y = 1.5
for title, desc, color in layers:
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(y), Inches(8.3), Inches(1.0))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.08), Inches(8), Inches(0.45))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(18)
    tp.font.bold = True
    tp.font.color.rgb = WHITE
    tp.alignment = PP_ALIGN.CENTER
    db = slide.shapes.add_textbox(Inches(2.7), Inches(y + 0.48), Inches(8), Inches(0.35))
    dp = db.text_frame.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(12)
    dp.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    dp.alignment = PP_ALIGN.CENTER
    if y < 4.5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.3), Inches(y + 1.0), Inches(0.5), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GRAY
        arrow.line.fill.background()
    y += 1.35

flow = slide.shapes.add_textbox(Inches(0.6), Inches(6.3), Inches(12), Inches(0.8))
fp = flow.text_frame.paragraphs[0]
fp.text = "데이터 흐름: 사용자 → Controller → Service → JdbcTemplate → DB  |  Python RPA → 크롤링 → 분석 → 차트"
fp.font.size = Pt(13)
fp.font.color.rgb = GRAY

# ── 슬라이드 5: 프론트엔드 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "프론트엔드", "Frontend Technology")

add_table(slide, ["기술", "설명"], [
    ["HTML5 / CSS3 / JavaScript", "시맨틱 마크업 · 스타일링 · 동적 UI"],
    ["Bootstrap 5", "반응형 그리드 · Navbar · Card · Modal"],
    ["Flexbox Layout", "헤더 · 카드 그리드 · 필터 바 · Footer"],
])

add_table(slide, ["화면", "기능"], [
    ["메인", "추천 매물 · 검색 바 · 카테고리"],
    ["매물 목록", "카드형 목록 · 페이지네이션 · 정렬"],
    ["매물 상세", "이미지 · 가격 · 면적 · 위치 정보"],
    ["로그인/회원가입", "Spring Security 연동 · BCrypt"],
    ["분석 대시보드", "지역별 시세 · Plotly 차트"],
], top=3.8)

# ── 슬라이드 6: 백엔드 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "백엔드", "Backend Technology")

add_table(slide, ["구분", "기술"], [
    ["Runtime", "JDK 21"],
    ["Framework", "Spring Boot 3.5.x"],
    ["View", "Thymeleaf (SSR)"],
    ["WAS", "Tomcat 10.1 (내장)"],
    ["Database", "MySQL 8"],
    ["Build", "Gradle"],
], top=1.4)

add_bullets(slide, [
    "Spring 의존성 (DI)",
    ("spring-boot-starter-web", 1),
    ("spring-boot-starter-security + BCrypt", 1),
    ("spring-boot-starter-thymeleaf", 1),
    ("mysql-connector-j", 1),
    ("spring-boot-starter-jdbc + JdbcTemplate", 1),
    ("Lombok", 1),
    "",
    "API: GET /properties · GET /properties/{id} · POST /auth/register · GET /api/analysis/region",
], left=6.5, top=1.4, width=6.5, font_size=15)

# ── 슬라이드 7: RPA + AI ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "RPA + AI (Python)", "Data Automation & Analysis")

add_table(slide, ["패키지", "용도"], [
    ["Requests", "HTTP 요청 · API 호출"],
    ["BeautifulSoup", "HTML 파싱 · 매물 크롤링"],
    ["NumPy / Pandas", "데이터 전처리 · 집계 · 분석"],
    ["Matplotlib / Seaborn", "정적 차트 · 통계 시각화"],
    ["Plotly", "인터랙티브 차트 · 대시보드"],
    ["Scikit-learn", "회귀·분류 · 가격 예측 ML"],
], top=1.4)

add_bullets(slide, [
    "자동화·분석 파이프라인",
    ("① RPA: 매물 사이트 → BeautifulSoup 파싱 → DB 저장", 1),
    ("② Pandas: 결측치·중복 처리 · 지역·평당가 집계", 1),
    ("③ 시각화: 지역별·월별 가격 추이 차트", 1),
    ("④ ML: 면적·층·지역 → 예상 가격 예측 모델", 1),
], top=5.0, font_size=16)

# ── 슬라이드 8: 배포 & 인프라 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "배포 & 인프라", "Deployment & Infrastructure")

add_table(slide, ["구분", "기술 / 서비스"], [
    ["Database", "TiDB Cloud (MySQL 8 호환)"],
    ["App Server", "Render 또는 AWS EC2"],
    ["Build", "Gradle — ./gradlew bootJar"],
    ["Version Control", "GitHub"],
    ["Security", "BCrypt · HTTPS · 환경변수 분리"],
], top=1.4)

add_bullets(slide, [
    "배포 절차",
    ("1. GitHub Push → 2. Gradle JAR 빌드 → 3. Render/EC2 배포 → 4. TiDB 연결", 1),
    "",
    "운영 고려사항",
    ("DB 연결 풀 · 백업 · HTTPS 적용", 1),
    ("민감 정보 환경변수 관리 (DB URL, Password)", 1),
], top=4.5, font_size=16)

# ── 슬라이드 9: 주요 기능 & 데모 ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "주요 기능 & 데모", "Key Features & Demo")

features = [
    "✔  매물 검색 · 필터 · 정렬",
    "✔  회원가입 · 로그인 (BCrypt 암호화)",
    "✔  Thymeleaf SSR 페이지 렌더링",
    "✔  Python RPA 매물 자동 수집",
    "✔  Pandas / Plotly 시세 분석 차트",
    "✔  Bootstrap 5 반응형 UI",
]

for i, feat in enumerate(features):
    col = 0 if i < 3 else 1
    row = i if i < 3 else i - 3
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.5 + row * 0.7)
    fb = slide.shapes.add_textbox(x, y, Inches(5.8), Inches(0.55))
    fp = fb.text_frame.paragraphs[0]
    fp.text = feat
    fp.font.size = Pt(17)
    fp.font.color.rgb = DARK

demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.0), Inches(12), Inches(2.5))
demo_box.fill.solid()
demo_box.fill.fore_color.rgb = LIGHT_BG
demo_box.line.color.rgb = PRIMARY

demo_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.15), Inches(11), Inches(0.4))
dtp = demo_title.text_frame.paragraphs[0]
dtp.text = "데모 시나리오 (3분)"
dtp.font.size = Pt(18)
dtp.font.bold = True
dtp.font.color.rgb = PRIMARY

add_bullets(slide, [
    ("1. 메인 → '강남' 지역 검색", 1),
    ("2. 매물 상세 페이지 확인", 1),
    ("3. 로그인 → 관심 매물 저장", 1),
    ("4. 분석 대시보드 → 지역별 시세 차트", 1),
], left=1.0, top=4.7, width=11, height=1.8, font_size=16)

# ── 슬라이드 10: 결론 & Q&A ──
slide = prs.slides.add_slide(BLANK)
add_header_bar(slide, "결론 & Q&A", "Conclusion")

add_bullets(slide, [
    "성과",
    ("Spring Boot 3.5 + JDK 21 최신 스택 적용", 1),
    ("웹 · DB · 보안 · 데이터 · RPA 풀스택 통합", 1),
    ("TiDB + Render/EC2 클라우드 배포 경험", 1),
    "",
    "한계 & 개선 방향",
    ("실시간 지도 API (Kakao/Naver Map) 연동", 1),
    ("REST API + JWT 분리 아키텍처", 1),
    ("크롤링 스케줄링 · 에러 재시도 자동화", 1),
    ("ML 모델 정확도 고도화", 1),
], font_size=17)

thanks = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11), Inches(0.8))
thp = thanks.text_frame.paragraphs[0]
thp.text = "감사합니다  |  Q & A"
thp.font.size = Pt(28)
thp.font.bold = True
thp.font.color.rgb = PRIMARY
thp.alignment = PP_ALIGN.CENTER

# 저장
output_path = os.path.join(os.path.dirname(__file__), "부동산웹플랫폼_발표자료.pptx")
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")
